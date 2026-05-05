from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xDB)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bar(slide, top, color):
    shape = slide.shapes.add_shape(1, Inches(0), top, Inches(13.333), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tb


def add_bullets(slide, left, top, width, height, items, size=12, color=WHITE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)


def make_slide(bg=DARK_BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg)
    add_bar(slide, Inches(0), ACCENT_BLUE)
    add_bar(slide, Inches(7.42), ACCENT_BLUE)
    return slide


# ── TITLE SLIDE ──
slide = make_slide()
add_text(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
         "ITAI 2374 \u2014 Robot Operating Systems & Platforms",
         size=36, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(0.8),
         "Individual Portfolio", size=28, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
         "Willie Lee Breaux III", size=24, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.7), Inches(11.3), Inches(0.5),
         "Team Apex AI  |  Houston Community College  |  Spring 2026",
         size=16, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ── TABLE OF CONTENTS ──
slide = make_slide()
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7),
         "Table of Contents", size=30, bold=True)
toc = [
    "Module S01 \u2014 Team Formation: Apex AI",
    "Module A01 \u2014 Obtaining and Reviewing the Raspberry Pi Book",
    "Module A04 \u2014 Synthesia AI Video Agent",
    "Module A05 \u2014 NVIDIA GTC Registration and Keynote Review",
    "Module L04 \u2014 Raspberry Pi Hands-On Lab: Audio, Camera, and Sense HAT",
    "Module L05 \u2014 PiCar-X Robot Lab",
    "Module A06 \u2014 Universal Robots UR3e Cobot: Introduction",
    "Module MT  \u2014 Midterm Project: UR3e Waypoint Pick-and-Place",
]
for i, item in enumerate(toc):
    y = Inches(1.5) + Inches(i * 0.6)
    c = ACCENT_BLUE if i % 2 == 0 else ACCENT_GREEN
    add_text(slide, Inches(1.5), y, Inches(0.5), Inches(0.5), f"{i+1}.", size=18, color=c, bold=True)
    add_text(slide, Inches(2.1), y, Inches(9), Inches(0.5), item, size=18)

# ── MODULE DATA ──
modules = [
    {
        "title": "Module S01: Team Formation \u2014 Apex AI",
        "activities": [
            "Formed Team Apex AI with teammates Kaden Glover, Williane Yarro, and Cesar Noriega",
            'Established team motto: "Control the systems. Trust the Model"',
            'Selected team song: "Power" by Kanye West',
            "Identified team strengths: problem-solving, Python/C++/ROS programming, hardware-software integration, applied AI (CV, Deep Learning, Edge AI), responsible robotics practices",
        ],
        "results": [
            "Established a collaborative team identity and aligned on shared technical strengths for the semester",
        ],
        "terms": [
            ("ROS (Robot Operating System)", "An open-source framework for writing robot software, providing tools and libraries for building robot applications"),
            ("Edge AI", "Running AI algorithms locally on hardware devices rather than in the cloud, enabling real-time processing"),
            ("GPIO (General-Purpose Input/Output)", "Pins on a board like Raspberry Pi used to interface with external hardware components"),
        ],
        "citations": ["None cited"],
        "indiv": [],
        "coll": ["Team formation, motto, song, strengths identification (no individual attribution in report)"],
    },
    {
        "title": "Module A01: Obtaining and Reviewing the Raspberry Pi Book",
        "activities": [
            'Ordered "The Official Raspberry Pi Beginner\'s Guide 6th Edition" from Amazon on first day of class',
            "Reviewed the book to identify what was most interesting and what was learned about Raspberry Pi",
            "Each team member reported individual findings from the book",
        ],
        "results": [
            "Found it most interesting that a Raspberry Pi can become a mini supercomputer by clustering Pi 5s",
            'Noted that adding an LCD screen, microphone, camera, and antenna was "a real eye-opening moment"',
            "Learned that Raspberry Pi can be used to learn code, build robots, and pursue creative projects",
            "Team concluded Raspberry Pi is a foundation for computing, engineering, and AI skills",
        ],
        "terms": [
            ("Raspberry Pi", "A single-board computer providing an affordable, flexible computing platform"),
            ("Raspberry Pi OS", "The official operating system for the Raspberry Pi desktop environment"),
            ("Raspberry Pi Imager", "Software tool to write Raspberry Pi OS to a microSD card for booting"),
            ("Clustering", "Connecting multiple Pi units together for more powerful computing"),
            ("HAT (Hardware Attached on Top)", "Modules that integrate onto the Pi's GPIO header"),
            ("SSD (Solid-State Drive)", "Storage device to enhance Pi's speed and performance"),
            ("Active Cooler", "A fan to prevent overheating and ensure consistent performance"),
        ],
        "citations": [
            "Upton, E., & Halfacree, G. (2021). The Official Raspberry Pi Beginner's Guide (4th ed.). Raspberry Pi Press.",
            "https://www.raspberrypi.com/products/the-official-raspberry-pi-beginners-guide/",
        ],
        "indiv": ["Ordering the book, interest in clustering Pi 5s, LCD/microphone/camera/antenna capabilities, Pi for code/robots/creative projects"],
        "coll": ["Team introduction paragraph, conclusion, shared source citation"],
    },
    {
        "title": "Module A04: Synthesia AI Video Agent",
        "activities": [
            "Created AI video agents on Synthesia.io to showcase HCC student AI capabilities and demonstrate skills to Intel",
            "Created a prompt/script instructing the AI agent on what to produce",
            "Produced two videos: one for HCC and one for Intel",
        ],
        "results": [
            "Learned that with just a prompt/script, anyone can create educational marketing materials, postcards, resume videos, and more",
            "Concluded that Synthesia's capabilities are endless and relevant to any field",
        ],
        "terms": [
            ("Synthesia.io", "An AI video generation platform using AI-generated avatars and voices"),
            ("AI Video Agent", "An AI system that generates video content autonomously from scripts/prompts"),
        ],
        "citations": [
            "HCC video: https://share.synthesia.io/ef59d3e7-a28a-4b1c-882e-f2f68689d9f4",
            "Intel video: https://share.synthesia.io/87b5658f-679d-4127-b787-64993feb5f35",
        ],
        "indiv": [],
        "coll": ["Entire report is unattributed team work (no individual sections)"],
    },
    {
        "title": "Module A05: NVIDIA GTC Registration and Keynote Review",
        "activities": [
            "Registered as an NVIDIA Developer (was unable to register for GTC itself)",
            "Used Developer access to view recorded GTC sessions after the event",
            "Viewed Jensen Huang's keynote address",
            "Selected sessions: Agentic AI, Robotics, Computer Vision, Edge Computing, and Simulation",
            "Watched the recorded session on Agentic AI",
        ],
        "results": [
            "Keynote made clear AI is a foundational system impacting nearly every industry",
            "Shift towards agentic AI and AI factories suggests greater automation",
            "Learned agentic AI can plan, make decisions, and complete tasks \u2014 not just respond to prompts",
            "Agentic AI works in a loop of observing, thinking, and acting",
            "Started implementing and integrating agentic AI into personal projects",
        ],
        "terms": [
            ("GTC (GPU Technology Conference)", "NVIDIA's annual AI/deep learning/accelerated computing conference"),
            ("NVIDIA Developer Program", "Access to NVIDIA tools, SDKs, recorded sessions, and resources"),
            ("Agentic AI", "AI that autonomously plans, decides, and acts in an observe-think-act loop"),
            ("AI Factory", "Computing infrastructure that converts data into intelligence via tokens"),
            ("CUDA", "NVIDIA's parallel computing platform for GPU-accelerated computing"),
            ("RTX", "NVIDIA's real-time ray tracing and AI-accelerated graphics platform"),
            ("DLSS (Deep Learning Super Sampling)", "AI rendering tech that boosts frames while generating sharp images"),
            ("CUDA-X", "NVIDIA GPU-accelerated libraries for healthcare, finance, robotics, telecom"),
            ("Vector Database", "Database for high-dimensional vector embeddings, used with AI models"),
            ("Edge Computing", "Processing data near the source for real-time AI applications"),
        ],
        "citations": [
            "NVIDIA GTC 2026 Keynote by Jensen Huang",
            "NVIDIA GTC recorded sessions (via NVIDIA Developer account)",
        ],
        "indiv": ["Registering as NVIDIA Developer, viewing Agentic AI session, keynote observations, session preferences, implementing agentic AI in personal projects"],
        "coll": ["Team introduction paragraph, conclusion"],
    },
    {
        "title": "Module L04: Raspberry Pi Hands-On Lab \u2014 Audio, Camera, Sense HAT",
        "activities": [
            "Built the Raspberry Pi case and inserted the SD card",
            "Wrote code to understand the Pi's foundation",
            "Connected speakers and a camera lens for hands-on hardware experience",
            "Worked with Sense HAT: pressure, humidity, temperature, color, orientation, movement sensors",
            "Wrote messages and changed LED matrix colors using Python in Thonny IDE",
            "Programmed Sense HAT accelerometer to detect movement and recall coordinates",
            "Debugged error caused by a missing value assignment using the '=' operator",
        ],
        "results": [
            "Handling the Raspberry Pi was exciting from start to finish",
            "Identified root cause of Thonny IDE error as a missing '=' operator",
            "Successfully demonstrated audio playback, camera output, and interactive Sense HAT functions",
            "Learned how hardware and software integrate in embedded systems",
            '"An extraordinary introduction to the capabilities of devices like the Raspberry Pi"',
        ],
        "terms": [
            ("Sense HAT", "Add-on board with sensors (pressure, humidity, temp, color, accelerometer) and 8x8 RGB LED matrix"),
            ("Thonny IDE", "Beginner-friendly Python IDE for writing/running scripts on the Raspberry Pi"),
            ("Accelerometer", "Sensor measuring acceleration forces and detecting movement/orientation changes"),
            ("GPIO Pins", "Physical pins for connecting/communicating with external hardware like the Sense HAT"),
            ("LED Matrix", "8x8 array of LEDs displaying text, colors, and patterns programmatically"),
            ("Embedded System", "A computer with a dedicated function within a larger system"),
        ],
        "citations": ["None explicitly cited in Willie's section"],
        "indiv": ["Building case, inserting SD card, writing foundational code, connecting speakers/camera, Sense HAT work, programming accelerometer, debugging missing '=' operator"],
        "coll": ["Team introduction, conclusion, lab setup activities described without individual attribution"],
    },
    {
        "title": "Module L05: PiCar-X Robot Lab",
        "activities": [
            "Assembled PiCar-X by connecting Raspberry Pi to the chassis",
            "Connected power supply, Ethernet cable, and hardware modules",
            "Installed PiCar-X software from SunFounder documentation",
            "Calibrated using python3 calibration.py and python3 cali_servo_motor.py",
            "Ran network test with ping google.com",
            "Navigated directories using cd picar-x, ls, pwd",
            "Resolved script naming error: video said 3.keyboard_control.py; correct was 2.keyboard_control.py",
            "Controlled robot via keyboard: forward, backward, turns, jump",
            "Missed first session; caught up with Cesar's help on a pre-built but non-operating robot",
        ],
        "results": [
            "Initially confusing due to missing the first session and returning to a non-functional robot",
            "With Cesar's help, caught up with the team and understood the system",
            "Learned patience, teamwork, and proper documentation make a big difference",
            "SunFounder docs were more reliable than the misleading instructional video",
            "After calibration and keyboard control activation, robot became functional",
            "Team successfully built, configured, calibrated, troubleshot, and controlled PiCar-X",
        ],
        "terms": [
            ("PiCar-X", "SunFounder robot car kit for Raspberry Pi with motors, servos, and programmable controls"),
            ("SunFounder", "Manufacturer of PiCar-X, providing docs and software for assembly/control"),
            ("Calibration (robotics)", "Adjusting motors/servos for accurate alignment and precise movement"),
            ("Servo Motor", "Motor allowing precise angular position control, used for PiCar-X steering"),
            ("ls command", "Lists files/directories; used to verify the correct script filename"),
            ("pwd command", "Prints current working directory for navigation verification"),
            ("ping command", "Network utility testing connectivity (e.g., ping google.com)"),
        ],
        "citations": ["SunFounder PiCar-X documentation (primary source for calibration and activation codes)"],
        "indiv": ["Missing first session, catching up with Cesar's help, initial confusion, learning patience/teamwork/documentation value, minor hardware fixes"],
        "coll": ["Assembly, software installation, calibration, script error resolution, directory navigation, network testing, robot control, conclusion"],
    },
    {
        "title": "Module A06: Universal Robots UR3e Cobot \u2014 Introduction",
        "activities": [
            "Reviewed basics of powering UR3e on/off via control box and PolyScope on teach pendant",
            "Studied safety protocols including E-Stop on control box and pendant",
            "Watched UR instructional videos prior to hands-on session",
            "Used teach pendant arrow controls for individual joint and orientation mode movement",
            "Explored the arm's range of motion and limits (unreachable positions, stopping angles)",
            "Positioned gripper over plastic cup, lowered, picked up, and relocated it",
            "Grabbed wooden block and dropped it into cup \u2014 multiple adjustments to align release point",
        ],
        "results": [
            "Gained familiarity with PolyScope interface and teach pendant controls",
            "Understood that even basic movements require high precision in positioning",
            "Learned difference between joint-by-joint control and orientation-based movement",
            "Successfully completed pick-and-place exercise (cup and block-into-cup)",
        ],
        "terms": [
            ("UR3e", "Collaborative robot arm by Universal Robots for light-duty industrial tasks"),
            ("Cobot (Collaborative Robot)", "Robot designed to work alongside humans with safety features"),
            ("PolyScope", "GUI on the UR teach pendant for programming and controlling the robot"),
            ("Teach Pendant", "Handheld tablet-like device to manually control and program the UR3e"),
            ("E-Stop (Emergency Stop)", "Safety button that immediately halts all robot movement"),
            ("Jogging (robotics)", "Manually moving a robot arm by individual joint or directional orientation"),
        ],
        "citations": ["Universal Robots UR instructional videos (viewed prior to hands-on session)"],
        "indiv": [],
        "coll": ["Entire report uses first-person plural with no individual attribution \u2014 all team-level"],
    },
    {
        "title": "Module MT: Midterm \u2014 UR3e Waypoint Pick-and-Place",
        "activities": [
            "Programmed UR3e to autonomously pick up blocks and transfer to conveyor belt via waypoints",
            "Configured 13 total waypoints using Move button and Set Waypoint in program interface",
            "Used MoveL for straight-line precision; MoveP for smooth transitions; Circle Move for curves",
            "Set tool speed 250 mm/s, TCP speed 250 mm/s, base acceleration 1200 mm/s\u00B2",
            "Resolved braking system issue \u2014 brakes hindering movement had to be disabled",
            "Fixed gripper release malfunction by adding an extra waypoint for block retrieval",
            "Completed all available UR training courses (except two unavailable)",
            "First day: explored Cobot movements, manually simulated conveyor-to-cup pick-and-place",
            "Resolved gripper issue with help from another team",
        ],
        "results": [
            "Robot autonomously moved 6 blocks onto conveyor belt without further human intervention",
            "MoveL was particularly beneficial for accurate pick-and-place",
            "Proper speed/acceleration configuration contributed to smooth, safe movement",
            "Learned Cobots are powerful when humans understand how to guide, program, and work alongside them",
            "Professors' insight: collaborative robots assist humans, not replace them",
            "Completed UR training courses for clearer understanding of Cobot operations",
            "Time with Cobot was cut short but provided wider understanding of AI/robotics in manufacturing",
        ],
        "terms": [
            ("Waypoint", "A programmed position in the robot's workspace defining a point in its motion path"),
            ("MoveL (Linear Move)", "Straight-line motion between waypoints for precision tasks"),
            ("MoveP (Process Move)", "Smoother path transitions between waypoints, blending movements"),
            ("Circle Move", "Motion command for curved/circular paths"),
            ("TCP (Tool Center Point)", "Reference point at tip of robot's tool for defining positions"),
            ("Tool Speed", "Velocity of the end effector during operation (250 mm/s)"),
            ("Base Acceleration", "Rate of velocity change (1200 mm/s\u00B2)"),
            ("Pick-and-Place", "Automated task: pick from one location, place at another"),
            ("Gripper", "End-of-arm tool for grasping and releasing objects"),
            ("Free-Drive Mode", "Human physically guides the robot arm for intuitive positioning"),
        ],
        "citations": [
            "Module 4: https://academy.universal-robots.com/free-e-learning/e-series-e-learning/e-series-core-track",
            "Universal Robots UR3e documentation and training courses",
        ],
        "indiv": ["Exploring Cobot movements, simulating pick-and-place, resolving gripper issue, completing UR training courses, observation that Cobots assist not replace humans"],
        "coll": ["13-waypoint programming, MoveL/MoveP/Circle Move config, speed/acceleration settings, braking fix, gripper malfunction fix, 6-block results, conclusion"],
    },
]

# Generate module slides
for mod in modules:
    # Title slide
    slide = make_slide()
    add_text(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
             mod["title"], size=32, bold=True, align=PP_ALIGN.CENTER)
    shape = slide.shapes.add_shape(1, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # Activities
    slide = make_slide()
    add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
             "Activities", size=26, color=ACCENT_BLUE, bold=True)
    items = ["\u2022  " + a for a in mod["activities"]]
    fs = 13 if len(items) <= 6 else 11 if len(items) <= 8 else 10
    add_bullets(slide, Inches(1.0), Inches(1.1), Inches(11.3), Inches(5.8), items, size=fs)

    # Results
    slide = make_slide()
    add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
             "Results & Learnings", size=26, color=ACCENT_GREEN, bold=True)
    items = ["\u2022  " + r for r in mod["results"]]
    fs = 13 if len(items) <= 5 else 11 if len(items) <= 7 else 10
    add_bullets(slide, Inches(1.0), Inches(1.1), Inches(11.3), Inches(5.8), items, size=fs)

    # Terms (paginated at 6 per slide)
    terms = mod["terms"]
    chunks = [terms[i:i+6] for i in range(0, len(terms), 6)]
    for chunk in chunks:
        slide = make_slide()
        add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
                 "Terms Learned", size=26, color=ACCENT_ORANGE, bold=True)
        y = Inches(1.1)
        for name, defn in chunk:
            add_text(slide, Inches(1.0), y, Inches(11.3), Inches(0.35),
                     name, size=13, color=ACCENT_ORANGE, bold=True)
            y += Inches(0.35)
            add_text(slide, Inches(1.3), y, Inches(11.0), Inches(0.55),
                     defn, size=11, color=LIGHT_GRAY)
            y += Inches(0.63)

    # Citations
    slide = make_slide()
    add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
             "Citations & References", size=26, color=ACCENT_PURPLE, bold=True)
    items = ["\u2022  " + c for c in mod["citations"]]
    add_bullets(slide, Inches(1.0), Inches(1.2), Inches(11.3), Inches(3.5), items, size=12, color=LIGHT_GRAY)

    # Contribution
    slide = make_slide()
    add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
             "Contribution Type", size=26, bold=True)
    y = Inches(1.3)
    if mod["indiv"]:
        add_text(slide, Inches(1.0), y, Inches(11.3), Inches(0.4),
                 "Individual (Willie Breaux):", size=15, color=ACCENT_GREEN, bold=True)
        y += Inches(0.5)
        for item in mod["indiv"]:
            add_text(slide, Inches(1.3), y, Inches(11.0), Inches(0.6),
                     "\u2022  " + item, size=12)
            y += Inches(0.55)
    y += Inches(0.2)
    if mod["coll"]:
        add_text(slide, Inches(1.0), y, Inches(11.3), Inches(0.4),
                 "Collective (Team Apex AI):", size=15, color=ACCENT_BLUE, bold=True)
        y += Inches(0.5)
        for item in mod["coll"]:
            add_text(slide, Inches(1.3), y, Inches(11.0), Inches(0.6),
                     "\u2022  " + item, size=12)
            y += Inches(0.55)

# ── CLOSING SLIDE ──
slide = make_slide()
add_text(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
         "Thank You", size=40, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.6),
         "Willie Lee Breaux III", size=24, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
         "Team Apex AI  |  ITAI 2374  |  Spring 2026  |  Houston Community College",
         size=16, color=MED_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
         '"Control the systems. Trust the Model."',
         size=18, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

# Save
output = r"C:\Users\wbrea\OneDrive\Documents\A.I. Robotics\Spring-Summer-Fall 2026\ITAI 2374 ROBOT OPERATING SYS & PLATFORM\Willie_Breaux_Portfolio_ITAI2374.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
